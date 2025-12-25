#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Create interactive HTML presentation from v2.pptx
Extract content and create professional interactive presentation
"""

from pptx import Presentation
from pptx.util import Pt
import html

def extract_slide_content(slide):
    """Extract text content from a slide"""
    content = {
        'title': '',
        'body': [],
    }
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                # Determine if title based on font size
                is_title = False
                try:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.size and run.font.size > Pt(32):
                                is_title = True
                                break
                        if is_title:
                            break
                except:
                    pass
                
                if is_title and not content['title']:
                    content['title'] = text
                else:
                    # Split into paragraphs preserving structure
                    paragraphs = [p.strip() for p in text.split('\n') if p.strip()]
                    content['body'].extend(paragraphs)
    
    return content

def create_interactive_html(pptx_file, output_file='presentation_v2.html'):
    """Create interactive HTML presentation from PowerPoint file"""
    print(f"טוען מצגת: {pptx_file}")
    prs = Presentation(pptx_file)
    
    slides_data = []
    for slide_idx, slide in enumerate(prs.slides):
        content = extract_slide_content(slide)
        slides_data.append(content)
        print(f"מעבד שקופית {slide_idx + 1}...")
        if content['title']:
            print(f"  כותרת: {content['title'][:50]}...")
        if content['body']:
            print(f"  פסקאות: {len(content['body'])}")
    
    # Create HTML
    html_content = f"""<!DOCTYPE html>
<html lang="he" dir="rtl">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>מצגת אינטראקטיבית</title>
    <style>
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}
        
        body {{
            font-family: 'Arial Hebrew', 'David', 'Gisha', 'Miriam', Arial, sans-serif;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: #141414;
            direction: rtl;
            overflow: hidden;
            height: 100vh;
        }}
        
        .presentation-wrapper {{
            width: 100%;
            height: 100vh;
            position: relative;
            display: flex;
            align-items: center;
            justify-content: center;
        }}
        
        .presentation-container {{
            width: 95%;
            height: 90vh;
            background: rgba(255, 255, 255, 0.98);
            border-radius: 20px;
            box-shadow: 0 20px 60px rgba(0, 0, 0, 0.3);
            overflow: hidden;
            position: relative;
            display: flex;
            flex-direction: column;
        }}
        
        .slide-container {{
            flex: 1;
            display: flex;
            align-items: center;
            justify-content: center;
            padding: 60px;
            overflow-y: auto;
            position: relative;
        }}
        
        .slide {{
            max-width: 1000px;
            width: 100%;
            padding: 50px;
            background: linear-gradient(135deg, #ffffff 0%, #f8f9fa 100%);
            border-radius: 15px;
            box-shadow: 0 10px 40px rgba(0, 0, 0, 0.1);
            display: none;
            opacity: 0;
            transform: scale(0.95);
            transition: all 0.4s cubic-bezier(0.4, 0, 0.2, 1);
        }}
        
        .slide.active {{
            display: block;
            opacity: 1;
            transform: scale(1);
            animation: slideIn 0.5s ease-out;
        }}
        
        @keyframes slideIn {{
            0% {{
                opacity: 0;
                transform: translateX(50px) scale(0.95);
            }}
            100% {{
                opacity: 1;
                transform: translateX(0) scale(1);
            }}
        }}
        
        .slide-title {{
            font-size: 52px;
            font-weight: bold;
            background: linear-gradient(135deg, #003366 0%, #0066CC 100%);
            -webkit-background-clip: text;
            -webkit-text-fill-color: transparent;
            background-clip: text;
            margin-bottom: 40px;
            padding-bottom: 25px;
            border-bottom: 4px solid #0066CC;
            text-align: right;
            line-height: 1.3;
            position: relative;
        }}
        
        .slide-title::after {{
            content: '';
            position: absolute;
            bottom: -4px;
            right: 0;
            width: 100px;
            height: 4px;
            background: linear-gradient(90deg, #0066CC, transparent);
        }}
        
        .slide-body {{
            font-size: 26px;
            color: #2c3e50;
            line-height: 2;
            text-align: right;
        }}
        
        .slide-body p {{
            margin-bottom: 25px;
            padding-right: 20px;
            border-right: 3px solid transparent;
            transition: all 0.3s ease;
        }}
        
        .slide-body p:hover {{
            border-right-color: #0066CC;
            padding-right: 30px;
        }}
        
        .controls {{
            position: fixed;
            bottom: 40px;
            left: 50%;
            transform: translateX(-50%);
            display: flex;
            gap: 20px;
            z-index: 1000;
            background: rgba(255, 255, 255, 0.95);
            backdrop-filter: blur(10px);
            padding: 20px 40px;
            border-radius: 50px;
            box-shadow: 0 8px 30px rgba(0, 0, 0, 0.2);
        }}
        
        .btn {{
            background: linear-gradient(135deg, #003366 0%, #0066CC 100%);
            color: white;
            border: none;
            padding: 15px 30px;
            font-size: 18px;
            font-family: 'Arial Hebrew', Arial, sans-serif;
            border-radius: 30px;
            cursor: pointer;
            transition: all 0.3s ease;
            font-weight: bold;
            box-shadow: 0 4px 15px rgba(0, 51, 102, 0.3);
        }}
        
        .btn:hover {{
            transform: translateY(-3px);
            box-shadow: 0 6px 20px rgba(0, 102, 204, 0.4);
        }}
        
        .btn:disabled {{
            background: #cccccc;
            cursor: not-allowed;
            transform: none;
        }}
        
        .slide-counter {{
            position: fixed;
            top: 30px;
            left: 30px;
            background: linear-gradient(135deg, #003366 0%, #0066CC 100%);
            color: white;
            padding: 15px 25px;
            border-radius: 30px;
            font-size: 18px;
            font-weight: bold;
            z-index: 1000;
            box-shadow: 0 4px 15px rgba(0, 51, 102, 0.3);
        }}
        
        .progress-bar {{
            position: fixed;
            top: 0;
            left: 0;
            height: 5px;
            background: linear-gradient(90deg, #0066CC, #003366);
            transition: width 0.4s ease;
            z-index: 1001;
            box-shadow: 0 2px 10px rgba(0, 102, 204, 0.5);
        }}
        
        .slide-thumbnails {{
            position: fixed;
            top: 50%;
            right: 20px;
            transform: translateY(-50%);
            display: flex;
            flex-direction: column;
            gap: 10px;
            z-index: 999;
            max-height: 70vh;
            overflow-y: auto;
            padding: 10px;
            background: rgba(255, 255, 255, 0.9);
            backdrop-filter: blur(10px);
            border-radius: 15px;
            box-shadow: 0 4px 20px rgba(0, 0, 0, 0.2);
        }}
        
        .thumbnail {{
            width: 80px;
            height: 60px;
            background: #f0f0f0;
            border-radius: 8px;
            cursor: pointer;
            transition: all 0.3s ease;
            border: 2px solid transparent;
            display: flex;
            align-items: center;
            justify-content: center;
            font-size: 12px;
            color: #666;
        }}
        
        .thumbnail:hover {{
            transform: scale(1.1);
            border-color: #0066CC;
        }}
        
        .thumbnail.active {{
            background: linear-gradient(135deg, #003366, #0066CC);
            color: white;
        }}
        
        @media (max-width: 768px) {{
            .slide {{
                padding: 30px;
            }}
            
            .slide-title {{
                font-size: 36px;
            }}
            
            .slide-body {{
                font-size: 20px;
            }}
            
            .slide-thumbnails {{
                display: none;
            }}
        }}
    </style>
</head>
<body>
    <div class="progress-bar" id="progressBar"></div>
    <div class="slide-counter" id="slideCounter"></div>
    
    <div class="slide-thumbnails" id="thumbnails"></div>
    
    <div class="presentation-wrapper">
        <div class="presentation-container">
            <div class="slide-container">
"""
    
    # Add slides
    for idx, slide_data in enumerate(slides_data):
        slide_class = "active" if idx == 0 else ""
        html_content += f'                <div class="slide {slide_class}" id="slide{idx}" data-index="{idx}">\n'
        
        if slide_data['title']:
            html_content += f'                    <h1 class="slide-title">{html.escape(slide_data["title"])}</h1>\n'
        
        if slide_data['body']:
            html_content += '                    <div class="slide-body">\n'
            for paragraph in slide_data['body']:
                html_content += f'                        <p>{html.escape(paragraph)}</p>\n'
            html_content += '                    </div>\n'
        else:
            html_content += '                    <div class="slide-body">\n'
            html_content += f'                        <p>שקופית {idx + 1}</p>\n'
            html_content += '                    </div>\n'
        
        html_content += '                </div>\n'
    
    html_content += """            </div>
        </div>
    </div>
    
    <div class="controls">
        <button class="btn" id="prevBtn" onclick="previousSlide()">← הקודם</button>
        <button class="btn" id="nextBtn" onclick="nextSlide()">הבא →</button>
    </div>
    
    <script>
        const slides = document.querySelectorAll('.slide');
        let currentSlide = 0;
        const totalSlides = slides.length;
        
        function updateSlide() {
            slides.forEach((slide, index) => {
                slide.classList.remove('active');
                if (index === currentSlide) {
                    setTimeout(() => {
                        slide.classList.add('active');
                    }, 50);
                }
            });
            
            document.getElementById('slideCounter').textContent = `${currentSlide + 1} / ${totalSlides}`;
            const progress = ((currentSlide + 1) / totalSlides) * 100;
            document.getElementById('progressBar').style.width = progress + '%';
            
            document.getElementById('prevBtn').disabled = currentSlide === 0;
            document.getElementById('nextBtn').disabled = currentSlide === totalSlides - 1;
            
            updateThumbnails();
        }
        
        function updateThumbnails() {
            const thumbnails = document.getElementById('thumbnails');
            thumbnails.innerHTML = '';
            
            slides.forEach((slide, index) => {
                const thumb = document.createElement('div');
                thumb.className = 'thumbnail' + (index === currentSlide ? ' active' : '');
                thumb.textContent = index + 1;
                thumb.onclick = () => goToSlide(index);
                thumbnails.appendChild(thumb);
            });
        }
        
        function goToSlide(index) {
            if (index >= 0 && index < totalSlides) {
                currentSlide = index;
                updateSlide();
            }
        }
        
        function nextSlide() {
            if (currentSlide < totalSlides - 1) {
                currentSlide++;
                updateSlide();
            }
        }
        
        function previousSlide() {
            if (currentSlide > 0) {
                currentSlide--;
                updateSlide();
            }
        }
        
        document.addEventListener('keydown', (e) => {
            if (e.key === 'ArrowRight' || e.key === 'ArrowDown' || e.key === ' ') {
                e.preventDefault();
                nextSlide();
            } else if (e.key === 'ArrowLeft' || e.key === 'ArrowUp') {
                e.preventDefault();
                previousSlide();
            } else if (e.key === 'Home') {
                goToSlide(0);
            } else if (e.key === 'End') {
                goToSlide(totalSlides - 1);
            }
        });
        
        let touchStartX = 0;
        let touchEndX = 0;
        
        document.addEventListener('touchstart', (e) => {
            touchStartX = e.changedTouches[0].screenX;
        });
        
        document.addEventListener('touchend', (e) => {
            touchEndX = e.changedTouches[0].screenX;
            const diff = touchStartX - touchEndX;
            if (Math.abs(diff) > 50) {
                if (diff > 0) {
                    nextSlide();
                } else {
                    previousSlide();
                }
            }
        });
        
        document.addEventListener('wheel', (e) => {
            if (e.deltaY > 0) {
                nextSlide();
            } else {
                previousSlide();
            }
        }, { passive: true });
        
        updateSlide();
    </script>
</body>
</html>"""
    
    # Write HTML file
    with open(output_file, 'w', encoding='utf-8') as f:
        f.write(html_content)
    
    print(f"\n✓ דף HTML אינטראקטיבי נוצר בהצלחה: {output_file}")
    print(f"  פתח את הקובץ בדפדפן כדי להציג את המצגת")

if __name__ == "__main__":
    input_file = "v2.pptx"
    output_file = "presentation_v2.html"
    
    try:
        create_interactive_html(input_file, output_file)
    except Exception as e:
        print(f"שגיאה: {e}")
        import traceback
        traceback.print_exc()


