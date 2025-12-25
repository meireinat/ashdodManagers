#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Create HTML presentation interface from PowerPoint file
- Extract slides content
- Create beautiful RTL HTML presentation
- Add navigation controls
- Optimized for daylight viewing
"""

from pptx import Presentation
from pptx.util import Pt
import html
import json
import os

def extract_slide_content(slide):
    """Extract text content from a slide"""
    content = {
        'title': '',
        'body': [],
        'shapes': []
    }
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                # Determine if title or body based on font size
                is_title = False
                try:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.size and run.font.size > Pt(32):
                                is_title = True
                                break
                        if is_title:
                            break
                except:
                    pass
                
                # Check shape name
                if shape.name and ('title' in shape.name.lower() or 'כותרת' in shape.name.lower()):
                    is_title = True
                
                if is_title and not content['title']:
                    content['title'] = text
                else:
                    # Split into paragraphs
                    paragraphs = [p.strip() for p in text.split('\n') if p.strip()]
                    content['body'].extend(paragraphs)
    
    return content

def create_html_presentation(pptx_file, output_file='presentation.html'):
    """Create HTML presentation from PowerPoint file"""
    print(f"טוען מצגת: {pptx_file}")
    prs = Presentation(pptx_file)
    
    slides_data = []
    for slide_idx, slide in enumerate(prs.slides):
        content = extract_slide_content(slide)
        slides_data.append(content)
        print(f"מעבד שקופית {slide_idx + 1}...")
    
    # Create HTML
    html_content = f"""<!DOCTYPE html>
<html lang="he" dir="rtl">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>ספריה דיגיטלית חכמה - מצגת</title>
    <style>
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}
        
        body {{
            font-family: 'Arial Hebrew', 'David', 'Gisha', 'Miriam', Arial, sans-serif;
            background: linear-gradient(135deg, #f5f8fa 0%, #ffffff 100%);
            color: #141414;
            direction: rtl;
            overflow: hidden;
            height: 100vh;
        }}
        
        .presentation-container {{
            width: 100%;
            height: 100vh;
            display: flex;
            flex-direction: column;
            position: relative;
        }}
        
        .slide-container {{
            flex: 1;
            display: flex;
            align-items: center;
            justify-content: center;
            padding: 40px;
            background: #ffffff;
            box-shadow: 0 4px 20px rgba(0, 0, 0, 0.1);
            margin: 20px;
            border-radius: 8px;
            overflow-y: auto;
        }}
        
        .slide {{
            max-width: 1200px;
            width: 100%;
            padding: 60px;
            background: #ffffff;
            border: 2px solid #dce6f0;
            border-radius: 12px;
            box-shadow: 0 8px 32px rgba(0, 51, 102, 0.1);
            display: none;
            animation: fadeIn 0.5s ease-in;
        }}
        
        .slide.active {{
            display: block;
        }}
        
        @keyframes fadeIn {{
            from {{
                opacity: 0;
                transform: translateX(20px);
            }}
            to {{
                opacity: 1;
                transform: translateX(0);
            }}
        }}
        
        .slide-title {{
            font-size: 48px;
            font-weight: bold;
            color: #003366;
            margin-bottom: 40px;
            padding-bottom: 20px;
            border-bottom: 3px solid #0066CC;
            text-align: right;
            line-height: 1.3;
        }}
        
        .slide-body {{
            font-size: 24px;
            color: #141414;
            line-height: 1.8;
            text-align: right;
        }}
        
        .slide-body p {{
            margin-bottom: 20px;
        }}
        
        .slide-body ul, .slide-body ol {{
            margin-right: 30px;
            margin-bottom: 20px;
        }}
        
        .slide-body li {{
            margin-bottom: 12px;
        }}
        
        .controls {{
            position: fixed;
            bottom: 30px;
            left: 50%;
            transform: translateX(-50%);
            display: flex;
            gap: 15px;
            z-index: 1000;
            background: rgba(255, 255, 255, 0.95);
            padding: 15px 30px;
            border-radius: 50px;
            box-shadow: 0 4px 20px rgba(0, 0, 0, 0.2);
        }}
        
        .btn {{
            background: #003366;
            color: white;
            border: none;
            padding: 12px 24px;
            font-size: 16px;
            font-family: 'Arial Hebrew', Arial, sans-serif;
            border-radius: 25px;
            cursor: pointer;
            transition: all 0.3s ease;
            font-weight: bold;
        }}
        
        .btn:hover {{
            background: #0066CC;
            transform: translateY(-2px);
            box-shadow: 0 4px 12px rgba(0, 102, 204, 0.4);
        }}
        
        .btn:active {{
            transform: translateY(0);
        }}
        
        .btn:disabled {{
            background: #cccccc;
            cursor: not-allowed;
            transform: none;
        }}
        
        .slide-counter {{
            position: fixed;
            top: 30px;
            left: 30px;
            background: rgba(0, 51, 102, 0.9);
            color: white;
            padding: 10px 20px;
            border-radius: 25px;
            font-size: 16px;
            font-weight: bold;
            z-index: 1000;
        }}
        
        .progress-bar {{
            position: fixed;
            top: 0;
            left: 0;
            height: 4px;
            background: #0066CC;
            transition: width 0.3s ease;
            z-index: 1001;
        }}
        
        .keyboard-hint {{
            position: fixed;
            bottom: 100px;
            left: 50%;
            transform: translateX(-50%);
            background: rgba(0, 0, 0, 0.7);
            color: white;
            padding: 10px 20px;
            border-radius: 20px;
            font-size: 14px;
            opacity: 0.7;
        }}
        
        @media (max-width: 768px) {{
            .slide {{
                padding: 30px;
            }}
            
            .slide-title {{
                font-size: 36px;
            }}
            
            .slide-body {{
                font-size: 20px;
            }}
            
            .controls {{
                bottom: 20px;
                padding: 10px 20px;
            }}
            
            .btn {{
                padding: 10px 18px;
                font-size: 14px;
            }}
        }}
    </style>
</head>
<body>
    <div class="progress-bar" id="progressBar"></div>
    <div class="slide-counter" id="slideCounter"></div>
    
    <div class="presentation-container">
        <div class="slide-container">
"""
    
    # Add slides
    for idx, slide_data in enumerate(slides_data):
        slide_class = "active" if idx == 0 else ""
        html_content += f'            <div class="slide {slide_class}" id="slide{idx}">\n'
        
        if slide_data['title']:
            html_content += f'                <h1 class="slide-title">{html.escape(slide_data["title"])}</h1>\n'
        
        if slide_data['body']:
            html_content += '                <div class="slide-body">\n'
            for paragraph in slide_data['body']:
                # Check if it's a list item
                if paragraph.startswith('•') or paragraph.startswith('-') or paragraph.startswith('*'):
                    html_content += f'                    <p>• {html.escape(paragraph.lstrip("•-* "))}</p>\n'
                else:
                    html_content += f'                    <p>{html.escape(paragraph)}</p>\n'
            html_content += '                </div>\n'
        
        html_content += '            </div>\n'
    
    html_content += """        </div>
    </div>
    
    <div class="controls">
        <button class="btn" id="prevBtn" onclick="previousSlide()">← הקודם</button>
        <button class="btn" id="nextBtn" onclick="nextSlide()">הבא →</button>
    </div>
    
    <div class="keyboard-hint">השתמש בחצים ← → לניווט</div>
    
    <script>
        const slides = document.querySelectorAll('.slide');
        let currentSlide = 0;
        const totalSlides = slides.length;
        
        function updateSlide() {
            slides.forEach((slide, index) => {
                slide.classList.remove('active');
                if (index === currentSlide) {
                    slide.classList.add('active');
                }
            });
            
            // Update counter
            document.getElementById('slideCounter').textContent = `${currentSlide + 1} / ${totalSlides}`;
            
            // Update progress bar
            const progress = ((currentSlide + 1) / totalSlides) * 100;
            document.getElementById('progressBar').style.width = progress + '%';
            
            // Update buttons
            document.getElementById('prevBtn').disabled = currentSlide === 0;
            document.getElementById('nextBtn').disabled = currentSlide === totalSlides - 1;
        }
        
        function nextSlide() {
            if (currentSlide < totalSlides - 1) {
                currentSlide++;
                updateSlide();
            }
        }
        
        function previousSlide() {
            if (currentSlide > 0) {
                currentSlide--;
                updateSlide();
            }
        }
        
        // Keyboard navigation
        document.addEventListener('keydown', (e) => {
            if (e.key === 'ArrowRight' || e.key === 'ArrowDown') {
                nextSlide();
            } else if (e.key === 'ArrowLeft' || e.key === 'ArrowUp') {
                previousSlide();
            } else if (e.key === 'Home') {
                currentSlide = 0;
                updateSlide();
            } else if (e.key === 'End') {
                currentSlide = totalSlides - 1;
                updateSlide();
            }
        });
        
        // Touch gestures for mobile
        let touchStartX = 0;
        let touchEndX = 0;
        
        document.addEventListener('touchstart', (e) => {
            touchStartX = e.changedTouches[0].screenX;
        });
        
        document.addEventListener('touchend', (e) => {
            touchEndX = e.changedTouches[0].screenX;
            handleSwipe();
        });
        
        function handleSwipe() {
            const swipeThreshold = 50;
            const diff = touchStartX - touchEndX;
            
            if (Math.abs(diff) > swipeThreshold) {
                if (diff > 0) {
                    // Swipe right (RTL) = next
                    nextSlide();
                } else {
                    // Swipe left (RTL) = previous
                    previousSlide();
                }
            }
        }
        
        // Initialize
        updateSlide();
        
        // Hide keyboard hint after 5 seconds
        setTimeout(() => {
            document.querySelector('.keyboard-hint').style.opacity = '0';
            setTimeout(() => {
                document.querySelector('.keyboard-hint').style.display = 'none';
            }, 500);
        }, 5000);
    </script>
</body>
</html>"""
    
    # Write HTML file
    with open(output_file, 'w', encoding='utf-8') as f:
        f.write(html_content)
    
    print(f"\n✓ ממשק HTML נוצר בהצלחה: {output_file}")
    print(f"  פתח את הקובץ בדפדפן כדי להציג את המצגת")

if __name__ == "__main__":
    input_file = "ספריה דיגיטלית חכמה.pptx"
    output_file = "presentation.html"
    
    try:
        create_html_presentation(input_file, output_file)
    except Exception as e:
        print(f"שגיאה: {e}")
        import traceback
        traceback.print_exc()

