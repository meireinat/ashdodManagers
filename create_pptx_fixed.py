#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Create PowerPoint presentation from HTML content - Fixed version
"""

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
import re
import html

def get_ashdod_port_colors():
    """Return Ashdod Port color scheme"""
    return {
        'primary': RGBColor(0, 51, 102),
        'secondary': RGBColor(0, 102, 204),
        'accent': RGBColor(255, 153, 0),
        'background': RGBColor(255, 255, 255),
        'text': RGBColor(20, 20, 20),
        'white': RGBColor(255, 255, 255),
    }

def parse_html_slides(html_file):
    """Parse HTML file and extract slide content"""
    with open(html_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    slides = []
    
    # Find all slide divs - improved regex
    slide_pattern = r'<div class="slide[^"]*" id="slide(\d+)">(.*?)</div>\s*</div>'
    matches = re.finditer(slide_pattern, content, re.DOTALL)
    
    for match in matches:
        slide_num = int(match.group(1))
        slide_content = match.group(2)
        
        # Extract title
        title_match = re.search(r'<h1 class="slide-title">(.*?)</h1>', slide_content)
        title = ''
        if title_match:
            title = html.unescape(title_match.group(1).strip())
        
        # Extract body paragraphs - improved extraction
        paragraphs = []
        
        # First, extract title area to exclude it
        title_end_pos = 0
        if title_match:
            title_end_pos = slide_content.find('</h1>') + 5
        
        # Find slide-body div - use non-greedy match
        body_match = re.search(r'<div class="slide-body">(.*?)</div>\s*</div>', slide_content[title_end_pos:], re.DOTALL)
        if not body_match:
            # Try without closing div
            body_match = re.search(r'<div class="slide-body">(.*?)</div>', slide_content[title_end_pos:], re.DOTALL)
        
        if body_match:
            body_text = body_match.group(1)
            # Find all <p> tags in body
            p_matches = re.finditer(r'<p>(.*?)</p>', body_text, re.DOTALL)
            for p_match in p_matches:
                para_text = p_match.group(1).strip()
                # Remove any remaining HTML tags
                para_text = re.sub(r'<[^>]+>', '', para_text)
                # Decode HTML entities like &quot;
                para_text = html.unescape(para_text)
                # Clean up whitespace
                para_text = ' '.join(para_text.split())
                if para_text:
                    paragraphs.append(para_text)
        
        # Debug: if still no paragraphs, try finding all <p> tags
        if not paragraphs:
            # Look for any <p> tags after title
            remaining_content = slide_content[title_end_pos:]
            p_matches = re.finditer(r'<p>(.*?)</p>', remaining_content, re.DOTALL)
            for p_match in p_matches:
                para_text = p_match.group(1).strip()
                para_text = re.sub(r'<[^>]+>', '', para_text)
                para_text = html.unescape(para_text)
                para_text = ' '.join(para_text.split())
                if para_text:
                    paragraphs.append(para_text)
        
        slides.append({
            'number': slide_num,
            'title': title,
            'paragraphs': paragraphs
        })
    
    return slides

def create_slide(prs, slide_data, colors):
    """Create a slide with the given content"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['background']
    
    # Add title if exists
    if slide_data['title']:
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1.5)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_ANCHOR.TOP
        title_frame.margin_left = Inches(0.2)
        title_frame.margin_right = Inches(0.2)
        title_frame.margin_top = Inches(0.1)
        title_frame.margin_bottom = Inches(0.1)
        
        # Clear and add paragraph
        title_frame.clear()
        p = title_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        p.space_after = Pt(0)
        
        # Set RTL
        try:
            pPr = p._element.get_or_add_pPr()
            pPr.set(qn('a:rtl'), '1')
        except:
            pass
        
        run = p.add_run()
        run.text = slide_data['title']
        
        # Set font with fallback
        try:
            run.font.name = 'Arial Hebrew'
        except:
            try:
                run.font.name = 'David'
            except:
                try:
                    run.font.name = 'Arial'
                except:
                    pass
        
        run.font.size = Pt(44)
        run.font.bold = True
        run.font.color.rgb = colors['primary']
        
        # Ensure text is visible
        if not run.text:
            run.text = slide_data['title']
    
    # Add body text
    if slide_data['paragraphs']:
        body_top = Inches(2.5) if slide_data['title'] else Inches(1)
        left = Inches(0.5)
        width = Inches(9)
        height = Inches(5.5)
        
        body_box = slide.shapes.add_textbox(left, body_top, width, height)
        body_frame = body_box.text_frame
        body_frame.word_wrap = True
        body_frame.vertical_anchor = MSO_ANCHOR.TOP
        body_frame.margin_left = Inches(0.3)
        body_frame.margin_right = Inches(0.3)
        body_frame.margin_top = Inches(0.2)
        body_frame.margin_bottom = Inches(0.2)
        
        # Clear default paragraph
        body_frame.clear()
        
        for idx, para_text in enumerate(slide_data['paragraphs']):
            p = body_frame.add_paragraph()
            p.alignment = PP_ALIGN.RIGHT
            p.space_after = Pt(12)
            p.line_spacing = 1.3
            
            # Set RTL
            try:
                pPr = p._element.get_or_add_pPr()
                pPr.set(qn('a:rtl'), '1')
            except:
                pass
            
            # Check if it's a title-like paragraph
            is_title_para = False
            emoji_starters = ('📚', '🤖', '💬', '📊', '📖', '🔍', '❌', '✅', '💡', '🎯', '📈')
            if any(para_text.startswith(emoji) for emoji in emoji_starters):
                is_title_para = True
            elif len(para_text) < 50 and ':' in para_text:
                is_title_para = True
            elif para_text.endswith('%'):
                is_title_para = True
            
            run = p.add_run()
            run.text = para_text
            
            # Set font with fallback
            try:
                run.font.name = 'Arial Hebrew'
            except:
                try:
                    run.font.name = 'David'
                except:
                    try:
                        run.font.name = 'Arial'
                    except:
                        pass
            
            if is_title_para:
                run.font.size = Pt(24)
                run.font.bold = True
                run.font.color.rgb = colors['primary']
            else:
                run.font.size = Pt(22)
                run.font.color.rgb = colors['text']
            
            # Ensure text is visible
            if not run.text:
                run.text = para_text

def create_presentation_from_html(html_file, output_file):
    """Create PowerPoint presentation from HTML file"""
    print(f"קורא קובץ HTML: {html_file}")
    slides_data = parse_html_slides(html_file)
    
    print(f"נמצאו {len(slides_data)} שקופיות")
    
    # Debug: print first slide content
    if slides_data:
        print(f"\nדוגמה - שקופית 1:")
        print(f"  כותרת: {slides_data[0]['title']}")
        print(f"  פסקאות: {len(slides_data[0]['paragraphs'])}")
        if slides_data[0]['paragraphs']:
            print(f"  פסקה ראשונה: {slides_data[0]['paragraphs'][0][:50]}...")
    
    # Create new presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    colors = get_ashdod_port_colors()
    
    # Sort slides by number
    slides_data.sort(key=lambda x: x['number'])
    
    # Create slides
    for slide_data in slides_data:
        print(f"יוצר שקופית {slide_data['number'] + 1}...")
        try:
            create_slide(prs, slide_data, colors)
        except Exception as e:
            print(f"  שגיאה בשקופית {slide_data['number'] + 1}: {e}")
            import traceback
            traceback.print_exc()
    
    print(f"\nשומר ל: {output_file}")
    prs.save(output_file)
    print("✓ המצגת נוצרה בהצלחה!")

if __name__ == "__main__":
    html_file = "presentation.html"
    output_file = "ספריה דיגיטלית חכמה - חדש.pptx"
    
    try:
        create_presentation_from_html(html_file, output_file)
    except Exception as e:
        print(f"שגיאה: {e}")
        import traceback
        traceback.print_exc()

