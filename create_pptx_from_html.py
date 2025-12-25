#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Create PowerPoint presentation from HTML content
Extract slides from presentation.html and create a new PPTX file
"""

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
import re
import html

def get_ashdod_port_colors():
    """Return Ashdod Port color scheme"""
    return {
        'primary': RGBColor(0, 51, 102),      # Deep navy blue
        'secondary': RGBColor(0, 102, 204),   # Ocean blue
        'accent': RGBColor(255, 153, 0),      # Port orange
        'background': RGBColor(255, 255, 255), # Pure white
        'text': RGBColor(20, 20, 20),         # Very dark
        'text_light': RGBColor(80, 80, 80),
        'white': RGBColor(255, 255, 255),
    }

def parse_html_slides(html_file):
    """Parse HTML file and extract slide content"""
    with open(html_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    slides = []
    
    # Find all slide divs
    slide_pattern = r'<div class="slide[^"]*" id="slide(\d+)">(.*?)</div>\s*</div>'
    matches = re.finditer(slide_pattern, content, re.DOTALL)
    
    for match in matches:
        slide_num = int(match.group(1))
        slide_content = match.group(2)
        
        # Extract title
        title_match = re.search(r'<h1 class="slide-title">(.*?)</h1>', slide_content)
        title = html.unescape(title_match.group(1)) if title_match else ''
        
        # Extract body paragraphs
        body_match = re.search(r'<div class="slide-body">(.*?)</div>', slide_content, re.DOTALL)
        body_text = body_match.group(1) if body_match else ''
        
        # Extract all paragraphs
        paragraphs = []
        p_matches = re.finditer(r'<p>(.*?)</p>', body_text, re.DOTALL)
        for p_match in p_matches:
            para_text = p_match.group(1).strip()
            # Remove HTML tags if any
            para_text = re.sub(r'<[^>]+>', '', para_text)
            # Decode HTML entities
            para_text = html.unescape(para_text)
            if para_text:
                paragraphs.append(para_text)
        
        slides.append({
            'number': slide_num,
            'title': title,
            'paragraphs': paragraphs
        })
    
    return slides

def set_rtl_direction(paragraph):
    """Set RTL direction for a paragraph using XML"""
    try:
        pPr = paragraph._element.get_or_add_pPr()
        
        # Set RTL using PowerPoint's namespace
        rtl_attr = qn('a:rtl')
        pPr.set(rtl_attr, '1')
        
        # Also try the text direction attribute
        try:
            pPr.set('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}dir', 'rtl')
            pPr.set('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}bidi', '1')
        except:
            pass
            
    except Exception as e:
        # Fallback: just set alignment
        paragraph.alignment = PP_ALIGN.RIGHT

def apply_hebrew_font(run, size=None, bold=False, color=None):
    """Apply Hebrew font with fallback options"""
    hebrew_fonts = ['Arial Hebrew', 'David', 'Gisha', 'Miriam', 'Arial', 'Calibri']
    
    font_set = False
    for font_name in hebrew_fonts:
        try:
            run.font.name = font_name
            font_set = True
            break
        except:
            continue
    
    # If no font was set, try setting a default
    if not font_set:
        try:
            run.font.name = 'Arial'
        except:
            pass
    
    if size:
        run.font.size = size
    if bold:
        run.font.bold = True
    if color:
        run.font.color.rgb = color
    
    # Ensure text is visible
    if not run.font.color.rgb:
        run.font.color.rgb = RGBColor(20, 20, 20)

def create_slide(prs, slide_data, colors):
    """Create a slide with the given content"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['background']
    
    # Add title if exists
    if slide_data['title']:
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1.5)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_ANCHOR.TOP
        title_frame.margin_left = Inches(0.2)
        title_frame.margin_right = Inches(0.2)
        
        # Clear default paragraph
        title_frame.clear()
        
        p = title_frame.add_paragraph()
        p.alignment = PP_ALIGN.RIGHT
        set_rtl_direction(p)
        
        run = p.add_run()
        run.text = slide_data['title']
        apply_hebrew_font(run, size=Pt(44), bold=True, color=colors['primary'])
    
    # Add body text
    if slide_data['paragraphs']:
        body_top = Inches(2.5) if slide_data['title'] else Inches(1)
        left = Inches(0.5)
        width = Inches(9)
        height = Inches(5.5)
        
        body_box = slide.shapes.add_textbox(left, body_top, width, height)
        body_frame = body_box.text_frame
        body_frame.word_wrap = True
        body_frame.vertical_anchor = MSO_ANCHOR.TOP
        body_frame.margin_left = Inches(0.3)
        body_frame.margin_right = Inches(0.3)
        
        # Clear default paragraph
        body_frame.clear()
        
        for idx, para_text in enumerate(slide_data['paragraphs']):
            p = body_frame.add_paragraph()
            
            p.alignment = PP_ALIGN.RIGHT
            p.space_after = Pt(12)
            p.line_spacing = 1.3
            set_rtl_direction(p)
            
            # Check if it's a title-like paragraph (starts with emoji or is short)
            is_title_para = False
            emoji_starters = ('📚', '🤖', '💬', '📊', '📖', '🔍', '❌', '✅', '💡', '🎯', '📈')
            if any(para_text.startswith(emoji) for emoji in emoji_starters):
                is_title_para = True
            elif len(para_text) < 50 and ':' in para_text:
                is_title_para = True
            elif para_text.endswith('%'):  # Percentage values
                is_title_para = True
            
            run = p.add_run()
            run.text = para_text
            
            if is_title_para:
                apply_hebrew_font(run, size=Pt(24), bold=True, color=colors['primary'])
            else:
                apply_hebrew_font(run, size=Pt(22), color=colors['text'])

def create_presentation_from_html(html_file, output_file):
    """Create PowerPoint presentation from HTML file"""
    print(f"קורא קובץ HTML: {html_file}")
    slides_data = parse_html_slides(html_file)
    
    print(f"נמצאו {len(slides_data)} שקופיות")
    
    # Create new presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    colors = get_ashdod_port_colors()
    
    # Sort slides by number
    slides_data.sort(key=lambda x: x['number'])
    
    # Create slides
    for slide_data in slides_data:
        print(f"יוצר שקופית {slide_data['number'] + 1}...")
        create_slide(prs, slide_data, colors)
    
    print(f"שומר ל: {output_file}")
    prs.save(output_file)
    print("\n✓ המצגת נוצרה בהצלחה!")

if __name__ == "__main__":
    html_file = "presentation.html"
    output_file = "ספריה דיגיטלית חכמה - חדש.pptx"
    
    try:
        create_presentation_from_html(html_file, output_file)
    except Exception as e:
        print(f"שגיאה: {e}")
        import traceback
        traceback.print_exc()

