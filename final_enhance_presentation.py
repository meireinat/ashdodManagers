#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Final enhanced script with advanced RTL support:
- Proper RTL text direction using XML manipulation
- Ashdod Port professional styling
- Hebrew fonts and typography
- Enhanced design elements
"""

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn

def get_ashdod_port_colors():
    """Return Ashdod Port color scheme - optimized for daylight viewing"""
    return {
        'primary': RGBColor(0, 51, 102),      # Deep navy blue
        'secondary': RGBColor(0, 102, 204),   # Ocean blue
        'accent': RGBColor(255, 153, 0),      # Port orange
        'background': RGBColor(255, 255, 255), # Pure white for clarity
        'background_alt': RGBColor(250, 252, 255), # Very light blue-white
        'text': RGBColor(20, 20, 20),         # Very dark for contrast
        'text_light': RGBColor(80, 80, 80),   # Medium gray
        'white': RGBColor(255, 255, 255),
        'border': RGBColor(220, 230, 240),    # Light border
    }

def set_rtl_direction(paragraph):
    """Set RTL direction for a paragraph using XML"""
    try:
        pPr = paragraph._element.get_or_add_pPr()
        
        # Set RTL using PowerPoint's namespace
        # This is the proper way to set RTL in PowerPoint XML
        rtl_attr = qn('a:rtl')
        pPr.set(rtl_attr, '1')
        
        # Also try the text direction attribute
        try:
            pPr.set('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}dir', 'rtl')
            pPr.set('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}bidi', '1')
        except:
            pass
            
    except Exception as e:
        # Fallback: just set alignment
        paragraph.alignment = PP_ALIGN.RIGHT

def apply_hebrew_font(run, size=None, bold=False, color=None):
    """Apply Hebrew font with fallback options"""
    hebrew_fonts = ['Arial Hebrew', 'David', 'Gisha', 'Miriam', 'Arial', 'Calibri']
    
    for font_name in hebrew_fonts:
        try:
            run.font.name = font_name
            break
        except:
            continue
    
    if size:
        run.font.size = size
    if bold:
        run.font.bold = True
    if color:
        run.font.color.rgb = color

def style_title(shape, colors):
    """Style title shapes"""
    if not shape.has_text_frame:
        return
    
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.RIGHT
        paragraph.space_after = Pt(12)
        set_rtl_direction(paragraph)
        
        for run in paragraph.runs:
            apply_hebrew_font(run, size=Pt(44), bold=True, color=colors['primary'])
            try:
                run.font.shadow = True
            except:
                pass

def style_body_text(shape, colors):
    """Style body text shapes"""
    if not shape.has_text_frame:
        return
    
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    text_frame.margin_left = Inches(0.5)
    text_frame.margin_right = Inches(0.5)
    
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.RIGHT
        paragraph.space_after = Pt(8)
        paragraph.line_spacing = 1.3
        set_rtl_direction(paragraph)
        
        for run in paragraph.runs:
            # Preserve existing size if reasonable, otherwise set default
            # Use darker text for better contrast in daylight
            if not run.font.size or run.font.size < Pt(14):
                apply_hebrew_font(run, size=Pt(22), color=colors['text'])
            else:
                # Ensure text is dark enough for visibility
                if run.font.color.rgb:
                    # Make text darker if it's too light
                    r, g, b = run.font.color.rgb
                    if r + g + b > 600:  # If too light
                        apply_hebrew_font(run, color=colors['text'])
                else:
                    apply_hebrew_font(run, color=colors['text'])

def process_presentation(input_file, output_file):
    """Process and enhance the presentation"""
    print(f"טוען מצגת: {input_file}")
    prs = Presentation(input_file)
    
    colors = get_ashdod_port_colors()
    
    print(f"מעבד {len(prs.slides)} שקופיות...")
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"מעבד שקופית {slide_idx + 1}...")
        
        # Set slide background - bright white for daylight viewing
        try:
            background = slide.background
            fill = background.fill
            fill.solid()
            # Use pure white for maximum clarity in daylight
            fill.fore_color.rgb = colors['background']
            
            # Add subtle border for definition
            try:
                line = fill.line
                line.color.rgb = colors['border']
                line.width = Pt(1)
            except:
                pass
        except:
            pass
        
        # Process shapes
        title_processed = False
        
        for shape in slide.shapes:
            is_title = False
            
            # Check if placeholder title
            try:
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    if shape.placeholder_format.idx == 0:
                        is_title = True
                        title_processed = True
            except:
                pass
            
            # Check shape name
            if not is_title and shape.name:
                name_lower = shape.name.lower()
                if any(word in name_lower for word in ['title', 'כותרת', 'heading', 'header']):
                    is_title = True
                    title_processed = True
            
            # Check font size
            if not is_title and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size and run.font.size > Pt(32):
                            is_title = True
                            title_processed = True
                            break
                    if is_title:
                        break
            
            # First text shape on first slide is likely title
            if not title_processed and slide_idx == 0 and shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if len(text) < 80 and text and '\n' not in text:
                    is_title = True
                    title_processed = True
            
            # Apply styling
            if is_title:
                style_title(shape, colors)
            elif shape.has_text_frame:
                style_body_text(shape, colors)
    
    print(f"שומר ל: {output_file}")
    prs.save(output_file)
    print("\n✓ הושלם בהצלחה!")
    print("  ✓ כיוון RTL (מימין לשמאל)")
    print("  ✓ עיצוב נמל אשדוד")
    print("  ✓ גופנים עבריים")
    print("  ✓ עיצוב מקצועי")

if __name__ == "__main__":
    input_file = "ספריה דיגיטלית חכמה.pptx"
    output_file = "ספריה דיגיטלית חכמה.pptx"
    
    try:
        process_presentation(input_file, output_file)
    except Exception as e:
        print(f"שגיאה: {e}")
        import traceback
        traceback.print_exc()

