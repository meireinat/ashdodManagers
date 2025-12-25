#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Fix "תנהלים" to "נהלים" in the first slide
"""

from pptx import Presentation

def fix_text_in_presentation(input_file, output_file):
    """Fix text in presentation"""
    print(f"טוען מצגת: {input_file}")
    prs = Presentation(input_file)
    
    # Fix only the first slide
    slide = prs.slides[0]
    changes_count = 0
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if 'תנהלים' in run.text:
                        run.text = run.text.replace('תנהלים', 'נהלים')
                        changes_count += 1
                        print(f"  עודכן: תנהלים → נהלים")
    
    print(f"שומר ל: {output_file}")
    prs.save(output_file)
    print(f"\n✓ הושלם! בוצעו {changes_count} שינויים")

if __name__ == "__main__":
    input_file = "ספריה דיגיטלית חכמה.pptx"
    output_file = "ספריה דיגיטלית חכמה.pptx"
    
    try:
        fix_text_in_presentation(input_file, output_file)
    except Exception as e:
        print(f"שגיאה: {e}")
        import traceback
        traceback.print_exc()


