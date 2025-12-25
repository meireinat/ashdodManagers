#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Read content from v1.pptx to see what's inside
"""

from pptx import Presentation
from pptx.util import Pt

def read_presentation(pptx_file):
    """Read and display presentation content"""
    print(f"קורא מצגת: {pptx_file}\n")
    prs = Presentation(pptx_file)
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"=== שקופית {slide_idx + 1} ===")
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    # Check font size
                    is_title = False
                    try:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.font.size and run.font.size > Pt(32):
                                    is_title = True
                                    print(f"[כותרת] {text[:100]}")
                                    break
                            if is_title:
                                break
                    except:
                        pass
                    
                    if not is_title:
                        print(f"[טקסט] {text[:100]}")
        
        print()

if __name__ == "__main__":
    read_presentation("v1.pptx")


