#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Update slide 2 title from "האתגר:" to "בעיות קיימות:"
"""

from pptx import Presentation

def update_slide2_title(input_file, output_file):
    """Update slide 2 title"""
    print(f"טוען מצגת: {input_file}")
    prs = Presentation(input_file)
    
    # Update slide 2 (index 1)
    slide = prs.slides[1]
    changes_count = 0
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if 'האתגר:' in run.text:
                        run.text = run.text.replace('האתגר:', 'בעיות קיימות:')
                        changes_count += 1
                        print(f"  עודכן: האתגר: → בעיות קיימות:")
    
    print(f"שומר ל: {output_file}")
    prs.save(output_file)
    print(f"\n✓ הושלם! בוצעו {changes_count} שינויים")

if __name__ == "__main__":
    input_file = "ספריה דיגיטלית חכמה.pptx"
    output_file = "ספריה דיגיטלית חכמה.pptx"
    
    try:
        update_slide2_title(input_file, output_file)
    except Exception as e:
        print(f"שגיאה: {e}")
        import traceback
        traceback.print_exc()


