#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Update text in PowerPoint presentation
Replace "מפוזר" with "מבוזר"
"""

from pptx import Presentation

def update_text_in_presentation(input_file, output_file):
    """Update text in presentation"""
    print(f"טוען מצגת: {input_file}")
    prs = Presentation(input_file)
    
    replacements = {
        'מפוזר': 'מבוזר',
        'מפוזרים': 'מבוזרים',
        'מפוזרות': 'מבוזרות',
    }
    
    changes_count = 0
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        original_text = run.text
                        new_text = original_text
                        
                        # Replace all occurrences
                        for old, new in replacements.items():
                            if old in new_text:
                                new_text = new_text.replace(old, new)
                                changes_count += new_text.count(new) - original_text.count(new)
                        
                        if new_text != original_text:
                            run.text = new_text
                            print(f"  שקופית {slide_idx + 1}: עודכן טקסט")
    
    print(f"שומר ל: {output_file}")
    prs.save(output_file)
    print(f"\n✓ הושלם! בוצעו {changes_count} שינויים")

if __name__ == "__main__":
    input_file = "ספריה דיגיטלית חכמה.pptx"
    output_file = "ספריה דיגיטלית חכמה.pptx"
    
    try:
        update_text_in_presentation(input_file, output_file)
    except Exception as e:
        print(f"שגיאה: {e}")
        import traceback
        traceback.print_exc()


